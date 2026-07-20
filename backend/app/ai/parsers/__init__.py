"""Document text parsers — no OCR (Module 6)."""

from __future__ import annotations

import csv
import io
import logging
import zipfile
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from xml.etree import ElementTree

logger = logging.getLogger(__name__)


@dataclass
class ExtractedPage:
    page: int
    text: str


@dataclass
class ExtractedDocument:
    text: str
    pages: list[ExtractedPage] = field(default_factory=list)
    parser: str = "unknown"


class BaseParser(ABC):
    @abstractmethod
    def extract(self, data: bytes, filename: str = "") -> ExtractedDocument:
        ...


class TxtParser(BaseParser):
    def extract(self, data: bytes, filename: str = "") -> ExtractedDocument:
        text = data.decode("utf-8", errors="replace")
        return ExtractedDocument(
            text=text, pages=[ExtractedPage(1, text)], parser="txt"
        )


class CsvParser(BaseParser):
    def extract(self, data: bytes, filename: str = "") -> ExtractedDocument:
        text = data.decode("utf-8", errors="replace")
        reader = csv.reader(io.StringIO(text))
        lines = [", ".join(row) for row in reader]
        joined = "\n".join(lines)
        return ExtractedDocument(
            text=joined, pages=[ExtractedPage(1, joined)], parser="csv"
        )


class PdfParser(BaseParser):
    def extract(self, data: bytes, filename: str = "") -> ExtractedDocument:
        try:
            from pypdf import PdfReader
        except ImportError as exc:
            raise RuntimeError(
                "pypdf is required for PDF parsing. pip install pypdf"
            ) from exc
        reader = PdfReader(io.BytesIO(data))
        pages: list[ExtractedPage] = []
        parts: list[str] = []
        for i, page in enumerate(reader.pages, start=1):
            t = (page.extract_text() or "").strip()
            pages.append(ExtractedPage(i, t))
            if t:
                parts.append(f"[Page {i}]\n{t}")
        return ExtractedDocument(
            text="\n\n".join(parts), pages=pages, parser="pdf"
        )


class DocxParser(BaseParser):
    def extract(self, data: bytes, filename: str = "") -> ExtractedDocument:
        try:
            from docx import Document as DocxDocument
        except ImportError:
            return self._fallback_zip_xml(data)
        doc = DocxDocument(io.BytesIO(data))
        paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        text = "\n".join(paras)
        return ExtractedDocument(
            text=text, pages=[ExtractedPage(1, text)], parser="docx"
        )

    def _fallback_zip_xml(self, data: bytes) -> ExtractedDocument:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            xml = zf.read("word/document.xml")
        root = ElementTree.fromstring(xml)
        ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        texts = [n.text for n in root.findall(".//w:t", ns) if n.text]
        text = "\n".join(texts)
        return ExtractedDocument(
            text=text, pages=[ExtractedPage(1, text)], parser="docx-xml"
        )


class PptxParser(BaseParser):
    def extract(self, data: bytes, filename: str = "") -> ExtractedDocument:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError(
                "python-pptx is required for PPTX parsing. pip install python-pptx"
            ) from exc
        prs = Presentation(io.BytesIO(data))
        pages: list[ExtractedPage] = []
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            chunks: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    chunks.append(shape.text.strip())
            slide_text = "\n".join(c for c in chunks if c)
            pages.append(ExtractedPage(i, slide_text))
            if slide_text:
                parts.append(f"[Slide {i}]\n{slide_text}")
        return ExtractedDocument(
            text="\n\n".join(parts), pages=pages, parser="pptx"
        )


class XlsxParser(BaseParser):
    def extract(self, data: bytes, filename: str = "") -> ExtractedDocument:
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise RuntimeError(
                "openpyxl is required for XLSX parsing. pip install openpyxl"
            ) from exc
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts: list[str] = []
        pages: list[ExtractedPage] = []
        for si, sheet in enumerate(wb.worksheets, start=1):
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(", ".join(cells))
            sheet_text = "\n".join(rows)
            pages.append(ExtractedPage(si, sheet_text))
            if sheet_text:
                parts.append(f"[Sheet {sheet.title}]\n{sheet_text}")
        return ExtractedDocument(
            text="\n\n".join(parts), pages=pages, parser="xlsx"
        )


class ParserFactory:
    """Return the best parser for a file extension."""

    _REGISTRY: dict[str, type[BaseParser]] = {
        "txt": TxtParser,
        "csv": CsvParser,
        "pdf": PdfParser,
        "docx": DocxParser,
        "pptx": PptxParser,
        "xlsx": XlsxParser,
    }

    @classmethod
    def get(cls, extension: str) -> BaseParser:
        ext = extension.lower().lstrip(".")
        parser_cls = cls._REGISTRY.get(ext)
        if parser_cls is None:
            raise ValueError(f"No text parser for extension: {ext}")
        return parser_cls()

    @classmethod
    def supported(cls) -> set[str]:
        return set(cls._REGISTRY.keys())
